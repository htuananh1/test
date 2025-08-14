import os
import re
import asyncio
import io
import logging
import random
import datetime
import base64
import json
from typing import Optional, Dict, List, Tuple, Any, Union
from collections import deque, defaultdict
from dataclasses import dataclass, field
from enum import Enum
import time
import aiohttp
from PIL import Image
import requests

from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.constants import ParseMode, ChatAction
from telegram.ext import (
    ApplicationBuilder, 
    ContextTypes, 
    MessageHandler, 
    filters, 
    CommandHandler, 
    CallbackQueryHandler
)
from telegram.error import BadRequest, TimedOut, NetworkError
from openai import OpenAI

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger("linh_bot")

# Constants for Vietnam features
VIETNAM_TIMEZONES = {
    "hanoi": "Asia/Ho_Chi_Minh",
    "saigon": "Asia/Ho_Chi_Minh", 
    "danang": "Asia/Ho_Chi_Minh"
}

VIETNAM_HOLIDAYS = {
    "01-01": "Tết Dương Lịch",
    "30-04": "Ngày Giải phóng miền Nam",
    "01-05": "Ngày Quốc tế Lao động",
    "02-09": "Ngày Quốc khánh"
}

VIETNAM_CITIES = {
    "hanoi": {"name": "Hà Nội", "lat": 21.0285, "lon": 105.8542},
    "hcm": {"name": "TP.HCM", "lat": 10.8231, "lon": 106.6297},
    "danang": {"name": "Đà Nẵng", "lat": 16.0544, "lon": 108.2022},
    "haiphong": {"name": "Hải Phòng", "lat": 20.8449, "lon": 106.6881},
    "cantho": {"name": "Cần Thơ", "lat": 10.0452, "lon": 105.7469},
    "nhatrang": {"name": "Nha Trang", "lat": 12.2388, "lon": 109.1967},
    "dalat": {"name": "Đà Lạt", "lat": 11.9404, "lon": 108.4583},
    "hue": {"name": "Huế", "lat": 16.4637, "lon": 107.5909}
}

@dataclass
class Config:
    BOT_TOKEN: str = field(default_factory=lambda: os.getenv("BOT_TOKEN", ""))
    VERCEL_API_KEY: str = field(default_factory=lambda: os.getenv("VERCEL_API_KEY", ""))
    BASE_URL: str = field(default_factory=lambda: os.getenv("BASE_URL", "https://ai-gateway.vercel.sh/v1"))
    
    # Text models (Vercel)
    CHAT_MODEL: str = field(default_factory=lambda: os.getenv("CHAT_MODEL", "anthropic/claude-3.5-haiku"))
    CODE_MODEL: str = field(default_factory=lambda: os.getenv("CODE_MODEL", "anthropic/claude-3.5-sonnet"))
    FILE_MODEL: str = field(default_factory=lambda: os.getenv("FILE_MODEL", "anthropic/claude-3.5-sonnet"))
    
    # Gemini models
    GEMINI_API_KEY: str = field(default_factory=lambda: os.getenv("GEMINI_API_KEY", ""))
    GEMINI_TEXT_MODEL: str = field(default_factory=lambda: os.getenv("GEMINI_TEXT_MODEL", "gemini-2.0-flash-exp"))
    GEMINI_VISION_MODEL: str = field(default_factory=lambda: os.getenv("GEMINI_VISION_MODEL", "gemini-1.5-flash"))
    GEMINI_IMAGE_GEN_MODEL: str = field(default_factory=lambda: os.getenv("GEMINI_IMAGE_GEN_MODEL", "gemini-2.0-flash-preview-image-generation"))
    
    MAX_TOKENS: int = field(default_factory=lambda: int(os.getenv("MAX_TOKENS", "1200")))
    MAX_TOKENS_CODE: int = field(default_factory=lambda: int(os.getenv("MAX_TOKENS_CODE", "4000")))
    FILE_OUTPUT_TOKENS: int = field(default_factory=lambda: int(os.getenv("FILE_OUTPUT_TOKENS", "6000")))
    
    CHUNK_CHARS: int = field(default_factory=lambda: int(os.getenv("CHUNK_CHARS", "120000")))
    PAGE_CHARS: int = field(default_factory=lambda: int(os.getenv("PAGE_CHARS", "3200")))
    CTX_TURNS: int = field(default_factory=lambda: int(os.getenv("CTX_TURNS", "15")))
    REQUEST_TIMEOUT: float = field(default_factory=lambda: float(os.getenv("REQUEST_TIMEOUT", "90")))
    
    WEATHER_API_KEY: str = field(default_factory=lambda: os.getenv("WEATHER_API_KEY", ""))
    NEWS_API_KEY: str = field(default_factory=lambda: os.getenv("NEWS_API_KEY", ""))
    
    CACHE_TTL: int = 3600
    MAX_CACHE_SIZE: int = 100

config = Config()

# Import Google AI libraries
try:
    import google.generativeai as genai
    GENAI_AVAILABLE = True
except ImportError:
    genai = None
    GENAI_AVAILABLE = False
    logger.warning("google-generativeai not installed")

# Import optional libraries
try:
    import chardet
except ImportError:
    chardet = None

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    import docx
    from docx import Document as DocxDocument
except ImportError:
    docx = None
    DocxDocument = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

class FileType(Enum):
    TEXT = "text"
    CODE = "code"
    PDF = "pdf"
    DOCX = "docx"
    XLSX = "xlsx"
    PPTX = "pptx"
    HTML = "html"
    JSON = "json"
    CSV = "csv"
    IMAGE = "image"
    ARCHIVE = "archive"
    UNKNOWN = "unknown"

ARCHIVES = (".zip", ".rar", ".7z", ".tar", ".tar.gz", ".tgz", ".tar.bz2", ".tar.xz")
TEXT_LIKE = (
    ".txt", ".md", ".log", ".csv", ".tsv", ".json", ".yaml", ".yml", 
    ".ini", ".cfg", ".env", ".xml", ".html", ".htm"
)
CODE_EXTENSIONS = (
    ".py", ".js", ".ts", ".java", ".c", ".cpp", ".cs", ".go", 
    ".php", ".rb", ".rs", ".sh", ".bat", ".ps1", ".sql", ".swift",
    ".kt", ".scala", ".r", ".m", ".dart", ".lua", ".pl", ".asm"
)
IMAGE_EXTENSIONS = (".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".svg", ".ico")

@dataclass
class ImageData:
    """Store image data for analysis"""
    file_id: str
    file_data: bytes
    caption: Optional[str] = None
    timestamp: float = field(default_factory=time.time)

@dataclass
class UserState:
    history: deque = field(default_factory=lambda: deque(maxlen=32))
    file_mode: bool = False
    image_mode: bool = False
    pending_file: Optional[Dict[str, Any]] = None
    pending_images: List[ImageData] = field(default_factory=list)
    last_result: str = ""
    active_messages: Dict[int, Any] = field(default_factory=dict)
    language: str = "vi"
    location: Optional[str] = "hanoi"

class BotState:
    def __init__(self):
        self.users: Dict[int, UserState] = defaultdict(UserState)
        self.pagers: Dict[Tuple[int, int], Dict] = {}
        self.cache: Dict[str, Tuple[Any, float]] = {}
        
    def get_user(self, chat_id: int) -> UserState:
        return self.users[chat_id]
    
    def cache_get(self, key: str) -> Optional[Any]:
        if key in self.cache:
            value, timestamp = self.cache[key]
            if time.time() - timestamp < config.CACHE_TTL:
                return value
            del self.cache[key]
        return None
    
    def cache_set(self, key: str, value: Any):
        if len(self.cache) >= config.MAX_CACHE_SIZE:
            oldest_key = min(self.cache.keys(), key=lambda k: self.cache[k][1])
            del self.cache[oldest_key]
        self.cache[key] = (value, time.time())

bot_state = BotState()

class VietnamServices:
    """Services related to Vietnam"""
    
    @staticmethod
    async def get_weather(city: str) -> str:
        """Get weather for Vietnamese cities"""
        if not config.WEATHER_API_KEY:
            return "❌ Thiếu API key thời tiết. Liên hệ @cucodoivandep"
        
        city_info = VIETNAM_CITIES.get(city.lower())
        if not city_info:
            cities = ", ".join(VIETNAM_CITIES.keys())
            return f"❌ Không tìm thấy {city}\nCác thành phố hỗ trợ: {cities}"
        
        try:
            url = f"https://api.openweathermap.org/data/2.5/weather"
            params = {
                "lat": city_info["lat"],
                "lon": city_info["lon"],
                "appid": config.WEATHER_API_KEY,
                "units": "metric",
                "lang": "vi"
            }
            
            async with aiohttp.ClientSession() as session:
                async with session.get(url, params=params) as response:
                    data = await response.json()
                    
            temp = data["main"]["temp"]
            feels_like = data["main"]["feels_like"]
            humidity = data["main"]["humidity"]
            description = data["weather"][0]["description"]
            wind_speed = data.get("wind", {}).get("speed", 0)
            
            weather_emoji = "☀️" if temp > 30 else "⛅" if temp > 20 else "🌧️"
            
            return (
                f"{weather_emoji} **Thời tiết {city_info['name']}**\n\n"
                f"🌡 Nhiệt độ: {temp}°C (cảm giác {feels_like}°C)\n"
                f"💧 Độ ẩm: {humidity}%\n"
                f"💨 Gió: {wind_speed} m/s\n"
                f"☁️ Trời: {description.capitalize()}\n\n"
                f"💡 Gợi ý: {'Nhớ mang ô ☂️' if 'mưa' in description.lower() else 'Thời tiết đẹp để đi chơi! 🌺'}"
            )
        except Exception as e:
            return f"❌ Lỗi lấy thời tiết: {str(e)[:100]}"
    
    @staticmethod
    async def get_exchange_rate() -> str:
        """Get USD/VND exchange rate"""
        try:
            url = "https://api.exchangerate-api.com/v4/latest/USD"
            async with aiohttp.ClientSession() as session:
                async with session.get(url) as response:
                    data = await response.json()
            
            vnd_rate = data["rates"].get("VND", 0)
            eur_rate = 1 / data["rates"].get("EUR", 1)
            gbp_rate = 1 / data["rates"].get("GBP", 1)
            jpy_rate = data["rates"].get("JPY", 0)
            cny_rate = data["rates"].get("CNY", 0)
            
            date = data.get("date", "")
            
            return (
                f"💱 **Tỷ giá hôm nay** ({date})\n\n"
                f"🇺🇸 1 USD = **{vnd_rate:,.0f}** VND\n"
                f"🇪🇺 1 EUR = **{vnd_rate * eur_rate:,.0f}** VND\n"
                f"🇬🇧 1 GBP = **{vnd_rate * gbp_rate:,.0f}** VND\n"
                f"🇯🇵 100 JPY = **{vnd_rate * 100 / jpy_rate:,.0f}** VND\n"
                f"🇨🇳 1 CNY = **{vnd_rate / cny_rate:,.0f}** VND\n\n"
                f"📈 Vàng SJC: ~92,000,000 VND/lượng"
            )
        except Exception as e:
            return f"❌ Lỗi lấy tỷ giá: {str(e)[:100]}"
    
    @staticmethod
    def get_vietnam_time() -> str:
        """Get current time in Vietnam"""
        try:
            import pytz
            vn_tz = pytz.timezone('Asia/Ho_Chi_Minh')
            vn_time = datetime.datetime.now(vn_tz)
            
            # Vietnamese day names
            vn_days = ['Thứ Hai', 'Thứ Ba', 'Thứ Tư', 'Thứ Năm', 'Thứ Sáu', 'Thứ Bảy', 'Chủ Nhật']
            day_name = vn_days[vn_time.weekday()]
            
            # Check if today is a holiday
            date_str = vn_time.strftime("%d-%m")
            holiday = VIETNAM_HOLIDAYS.get(date_str, "")
            
            # Calculate lunar calendar (approximate)
            lunar_info = "🌙 Âm lịch: Đang cập nhật"
            
            time_str = (
                f"🇻🇳 **Giờ Việt Nam**\n\n"
                f"📅 {day_name}, {vn_time.strftime('%d/%m/%Y')}\n"
                f"🕐 {vn_time.strftime('%H:%M:%S')} (GMT+7)\n"
                f"{lunar_info}"
            )
            
            if holiday:
                time_str += f"\n\n🎉 **{holiday}**"
            
            # Add greeting based on time
            hour = vn_time.hour
            if 5 <= hour < 11:
                greeting = "🌅 Chào buổi sáng!"
            elif 11 <= hour < 13:
                greeting = "☀️ Chào buổi trưa!"  
            elif 13 <= hour < 18:
                greeting = "🌤 Chào buổi chiều!"
            else:
                greeting = "🌙 Chào buổi tối!"
            
            time_str += f"\n\n{greeting}"
            
            return time_str
        except:
            vn_time = datetime.datetime.utcnow() + datetime.timedelta(hours=7)
            return f"🕐 Giờ VN: {vn_time.strftime('%H:%M:%S %d/%m/%Y')}"
    
    @staticmethod
    async def get_news() -> str:
        """Get Vietnamese news headlines"""
        if not config.NEWS_API_KEY:
            # Return some default news sources
            return (
                "📰 **Báo chí Việt Nam**\n\n"
                "🔸 VnExpress: vnexpress.net\n"
                "🔸 Tuổi Trẻ: tuoitre.vn\n"
                "🔸 Thanh Niên: thanhnien.vn\n"
                "🔸 Dân Trí: dantri.com.vn\n"
                "🔸 VTC News: vtc.vn\n\n"
                "💡 Cần API key để xem tin tức tự động"
            )
        
        try:
            url = "https://newsapi.org/v2/top-headlines"
            params = {
                "country": "vn",
                "apiKey": config.NEWS_API_KEY,
                "pageSize": 5
            }
            
            async with aiohttp.ClientSession() as session:
                async with session.get(url, params=params) as response:
                    data = await response.json()
            
            if data.get("status") != "ok":
                return "❌ Không lấy được tin tức"
            
            articles = data.get("articles", [])
            if not articles:
                return "📰 Không có tin tức mới"
            
            news_text = "📰 **Tin tức Việt Nam mới nhất**\n\n"
            for i, article in enumerate(articles[:5], 1):
                title = article.get("title", "")
                description = article.get("description", "")[:100]
                source = article.get("source", {}).get("name", "")
                news_text += f"**{i}. {title}**\n"
                if description:
                    news_text += f"_{description}_\n"
                if source:
                    news_text += f"📌 {source}\n"
                news_text += "\n"
            
            return news_text
        except Exception as e:
            return f"❌ Lỗi lấy tin tức: {str(e)[:100]}"

vietnam_services = VietnamServices()

class GeminiHandler:
    """Handle Gemini API for both vision and image generation"""
    
    def __init__(self):
        self.configured = False
        if config.GEMINI_API_KEY and GENAI_AVAILABLE:
            genai.configure(api_key=config.GEMINI_API_KEY)
            self.configured = True
    
    async def analyze_image(self, image_data: bytes, prompt: str) -> Optional[str]:
        """Analyze image using Gemini Vision"""
        if not self.configured:
            return None
        
        try:
            # Create model for vision
            model = genai.GenerativeModel(config.GEMINI_VISION_MODEL)
            
            # Convert bytes to PIL Image
            image = Image.open(io.BytesIO(image_data))
            
            # Generate content with image
            response = await asyncio.to_thread(
                model.generate_content,
                [prompt, image]
            )
            
            return response.text
            
        except Exception as e:
            logger.error(f"Gemini Vision error: {e}")
            return f"❌ Lỗi phân tích: {str(e)[:100]}"
    
    async def generate_image(self, prompt: str) -> Optional[bytes]:
        """Generate image using Gemini"""
        if not self.configured:
            return None
        
        try:
            # Use the image generation model
            model = genai.GenerativeModel(config.GEMINI_IMAGE_GEN_MODEL)
            
            # Generate image
            response = await asyncio.to_thread(
                model.generate_content,
                prompt
            )
            
            # Extract image data if available
            if hasattr(response, '_result') and hasattr(response._result, 'candidates'):
                for candidate in response._result.candidates:
                    if hasattr(candidate, 'content') and hasattr(candidate.content, 'parts'):
                        for part in candidate.content.parts:
                            if hasattr(part, 'inline_data') and part.inline_data:
                                return base64.b64decode(part.inline_data.data)
            
            return None
            
        except Exception as e:
            logger.error(f"Gemini Image Gen error: {e}")
            return None
    
    async def analyze_multiple_images(self, images: List[ImageData], prompt: str) -> str:
        """Analyze multiple images"""
        if not self.configured:
            return "❌ Gemini API chưa được cấu hình"
        
        results = []
        
        for i, img_data in enumerate(images, 1):
            vn_prompt = f"Ảnh {i}: {prompt}. Trả lời bằng tiếng Việt, chú ý các yếu tố văn hóa Việt Nam nếu có."
            
            result = await self.analyze_image(img_data.file_data, vn_prompt)
            
            if result:
                results.append(f"**📸 Ảnh {i}:**\n{result}")
            else:
                results.append(f"**📸 Ảnh {i}:** Không thể phân tích")
        
        return "\n\n".join(results) if results else "❌ Không thể phân tích ảnh"

gemini_handler = GeminiHandler()

class AIClient:
    """Handle text-based AI using Vercel API"""
    
    def __init__(self):
        self.client = None
        if config.VERCEL_API_KEY:
            self.client = OpenAI(
                api_key=config.VERCEL_API_KEY,
                base_url=config.BASE_URL,
                timeout=config.REQUEST_TIMEOUT
            )
    
    def stream_complete(
        self,
        model: str,
        messages: List[Dict[str, str]],
        max_tokens: int,
        temperature: float = 0.7
    ):
        if not self.client:
            yield "❌ Thiếu VERCEL_API_KEY."
            return
        
        try:
            stream = self.client.chat.completions.create(
                model=model,
                max_tokens=max_tokens,
                temperature=temperature,
                messages=messages,
                stream=True
            )
            
            for chunk in stream:
                if chunk.choices[0].delta.content:
                    yield chunk.choices[0].delta.content
                    
        except Exception as e:
            yield f"\n❌ Lỗi: {str(e)[:200]}"
    
    async def complete(
        self,
        model: str,
        messages: List[Dict[str, str]],
        max_tokens: int,
        temperature: float = 0.7,
        retries: int = 3
    ) -> str:
        if not self.client:
            return "❌ Thiếu VERCEL_API_KEY."
        
        cache_key = f"{model}:{json.dumps(messages)}:{max_tokens}:{temperature}"
        cached = bot_state.cache_get(cache_key)
        if cached:
            return cached
        
        last_error = None
        for attempt in range(retries):
            try:
                response = await asyncio.wait_for(
                    asyncio.to_thread(
                        self._sync_complete,
                        model, messages, max_tokens, temperature
                    ),
                    timeout=config.REQUEST_TIMEOUT + 10
                )
                
                bot_state.cache_set(cache_key, response)
                return response
                    
            except asyncio.TimeoutError:
                last_error = "Timeout - yêu cầu mất quá nhiều thời gian"
            except Exception as e:
                last_error = str(e)
                
            if attempt < retries - 1:
                await asyncio.sleep(1.5 * (attempt + 1) + random.random())
        
        raise Exception(f"Lỗi sau {retries} lần thử: {last_error}")
    
    def _sync_complete(
        self, 
        model: str, 
        messages: List[Dict[str, str]], 
        max_tokens: int,
        temperature: float
    ) -> str:
        response = self.client.chat.completions.create(
            model=model,
            max_tokens=max_tokens,
            temperature=temperature,
            messages=messages
        )
        return (response.choices[0].message.content or "").strip()

ai_client = AIClient()

class TextProcessor:
    @staticmethod
    def escape_markdown_v2(text: str) -> str:
        special_chars = ['_', '*', '[', ']', '(', ')', '~', '`', '>', '#', '+', '-', '=', '|', '{', '}', '.', '!']
        for char in special_chars:
            text = text.replace(char, f'\\{char}')
        return text
    
    @staticmethod
    def chunk_text(text: str, max_length: int = 4096) -> List[str]:
        """Split text into chunks for Telegram messages"""
        if len(text) <= max_length:
            return [text]
        
        chunks = []
        current = ""
        
        for line in text.split('\n'):
            if len(current) + len(line) + 1 > max_length:
                if current:
                    chunks.append(current)
                current = line
            else:
                current = current + '\n' + line if current else line
        
        if current:
            chunks.append(current)
        
        return chunks

text_processor = TextProcessor()

class FileProcessor:
    @staticmethod
    def detect_file_type(filename: str) -> FileType:
        name_lower = filename.lower()
        
        if any(name_lower.endswith(ext) for ext in ARCHIVES):
            return FileType.ARCHIVE
        elif any(name_lower.endswith(ext) for ext in CODE_EXTENSIONS):
            return FileType.CODE
        elif any(name_lower.endswith(ext) for ext in IMAGE_EXTENSIONS):
            return FileType.IMAGE
        elif name_lower.endswith('.pdf'):
            return FileType.PDF
        elif name_lower.endswith(('.docx', '.doc')):
            return FileType.DOCX
        elif name_lower.endswith(('.xlsx', '.xls')):
            return FileType.XLSX
        elif name_lower.endswith(('.pptx', '.ppt')):
            return FileType.PPTX
        elif name_lower.endswith(('.html', '.htm')):
            return FileType.HTML
        elif name_lower.endswith('.json'):
            return FileType.JSON
        elif name_lower.endswith('.csv'):
            return FileType.CSV
        elif any(name_lower.endswith(ext) for ext in TEXT_LIKE):
            return FileType.TEXT
        else:
            return FileType.UNKNOWN
    
    @staticmethod
    def detect_encoding(data: bytes) -> str:
        if not data:
            return "utf-8"
        
        if chardet:
            result = chardet.detect(data)
            if result and result.get('encoding'):
                return result['encoding']
        
        for encoding in ['utf-8', 'utf-16', 'latin-1', 'cp1252', 'gb2312']:
            try:
                data.decode(encoding)
                return encoding
            except UnicodeDecodeError:
                continue
        
        return 'utf-8'
    
    @staticmethod
    async def process_file(filename: str, data: bytes) -> Tuple[Optional[str], str]:
        file_type = FileProcessor.detect_file_type(filename)
        
        try:
            if file_type == FileType.ARCHIVE:
                return None, "📦 File nén - vui lòng giải nén trước khi gửi"
            
            elif file_type == FileType.IMAGE:
                return None, "image"
            
            elif file_type == FileType.PDF:
                if not PyPDF2:
                    return None, "❌ Cần cài đặt PyPDF2 để đọc PDF"
                return FileProcessor._read_pdf(data), "pdf"
            
            elif file_type == FileType.DOCX:
                if not docx:
                    return None, "❌ Cần cài đặt python-docx để đọc DOCX"
                return FileProcessor._read_docx(data), "docx"
            
            elif file_type == FileType.XLSX:
                if not openpyxl:
                    return None, "❌ Cần cài đặt openpyxl để đọc Excel"
                return FileProcessor._read_excel(data), "xlsx"
            
            elif file_type == FileType.PPTX:
                if not Presentation:
                    return None, "❌ Cần cài đặt python-pptx để đọc PowerPoint"
                return FileProcessor._read_pptx(data), "pptx"
            
            elif file_type == FileType.HTML:
                text = FileProcessor._decode_text(data)
                if BeautifulSoup:
                    soup = BeautifulSoup(text, 'html.parser')
                    return soup.get_text(separator='\n'), "html"
                return text, "html"
            
            elif file_type == FileType.JSON:
                text = FileProcessor._decode_text(data)
                try:
                    obj = json.loads(text)
                    return json.dumps(obj, indent=2, ensure_ascii=False), "json"
                except:
                    return text, "json"
            
            elif file_type in [FileType.TEXT, FileType.CODE, FileType.CSV]:
                return FileProcessor._decode_text(data), file_type.value
            
            else:
                return None, f"❌ File {filename} chưa được hỗ trợ"
                
        except Exception as e:
            return None, f"❌ Lỗi xử lý file: {str(e)[:100]}"
    
    @staticmethod
    def _decode_text(data: bytes) -> str:
        encoding = FileProcessor.detect_encoding(data)
        try:
            return data.decode(encoding, errors='ignore')
        except:
            return data.decode('utf-8', errors='ignore')
    
    @staticmethod
    def _read_pdf(data: bytes) -> str:
        pdf_file = io.BytesIO(data)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = []
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text.append(page.extract_text())
        return '\n'.join(text)
    
    @staticmethod
    def _read_docx(data: bytes) -> str:
        doc_file = io.BytesIO(data)
        doc = DocxDocument(doc_file)
        return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
    
    @staticmethod
    def _read_excel(data: bytes) -> str:
        excel_file = io.BytesIO(data)
        wb = openpyxl.load_workbook(excel_file, read_only=True)
        text = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text.append(f"=== Sheet: {sheet_name} ===")
            for row in sheet.iter_rows(values_only=True):
                row_text = '\t'.join([str(cell) if cell else '' for cell in row])
                if row_text.strip():
                    text.append(row_text)
        return '\n'.join(text)
    
    @staticmethod
    def _read_pptx(data: bytes) -> str:
        pptx_file = io.BytesIO(data)
        prs = Presentation(pptx_file)
        text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            text.append(f"=== Slide {slide_num} ===")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        text.append(shape.text)
        return '\n'.join(text)

file_processor = FileProcessor()

class MessageBuilder:
    @staticmethod
    def build_system_prompt(context_type: str = "chat") -> str:
        try:
            import pytz
            vn_tz = pytz.timezone('Asia/Ho_Chi_Minh')
            vn_time = datetime.datetime.now(vn_tz)
            time_str = vn_time.strftime('%d/%m/%Y %H:%M:%S')
        except:
            vn_time = datetime.datetime.utcnow() + datetime.timedelta(hours=7)
            time_str = vn_time.strftime('%d/%m/%Y %H:%M:%S')
        
        base_prompt = (
            f"Bạn là Linh - AI assistant thông minh của Việt Nam.\n"
            f"Thời gian Việt Nam: {time_str}\n"
            f"Được phát triển bởi Hoàng Tuấn Anh (@cucodoivandep)\n\n"
            f"Kiến thức sâu về:\n"
            f"• Lịch sử, văn hóa, địa lý Việt Nam\n"
            f"• Ẩm thực, du lịch, phong tục tập quán\n"
            f"• Tiếng Việt và các phương ngữ\n"
            f"• Kinh tế, xã hội Việt Nam\n\n"
        )
        
        if context_type == "chat":
            return base_prompt + (
                "Hướng dẫn giao tiếp:\n"
                "• Ưu tiên trả lời về Việt Nam khi phù hợp\n"
                "• Sử dụng tiếng Việt tự nhiên, thân thiện\n"
                "• Có thể dùng emoji phù hợp\n"
                "• Trả lời ngắn gọn, súc tích, đi thẳng vấn đề\n"
                "• Thể hiện sự am hiểu văn hóa Việt"
            )
        
        elif context_type == "code":
            return base_prompt + (
                "Bạn là lập trình viên Việt Nam chuyên nghiệp.\n\n"
                "Nguyên tắc:\n"
                "• Code sạch, dễ đọc, dễ bảo trì\n"
                "• Comment bằng tiếng Việt rõ ràng\n"
                "• Đặt tên biến/hàm theo chuẩn quốc tế\n"
                "• Xử lý tốt Unicode tiếng Việt\n"
                "• Tối ưu cho production"
            )
        
        elif context_type == "file":
            return base_prompt + (
                "Xử lý file với chuyên môn cao.\n\n"
                "Hướng dẫn:\n"
                "• Phân tích nội dung chính xác\n"
                "• Tóm tắt súc tích các điểm chính\n"
                "• Trích xuất thông tin quan trọng\n"
                "• Đề xuất cải thiện nếu phù hợp"
            )
        
        elif context_type == "image_context":
            return base_prompt + (
                "Bạn đang trả lời câu hỏi về hình ảnh.\n\n"
                "Lưu ý:\n"
                "• Dựa trên phân tích ảnh để trả lời\n"
                "• Chú ý chi tiết trong ảnh\n"
                "• Liên hệ với văn hóa Việt Nam nếu có\n"
                "• Trả lời chính xác, cụ thể"
            )
        
        return base_prompt
    
    @staticmethod
    def build_messages(
        chat_id: int,
        user_text: str,
        context_type: str = "chat",
        include_history: bool = True
    ) -> List[Dict[str, str]]:
        messages = []
        
        system_prompt = MessageBuilder.build_system_prompt(context_type)
        messages.append({"role": "system", "content": system_prompt})
        
        if include_history:
            user = bot_state.get_user(chat_id)
            history_messages = []
            
            keep_turns = config.CTX_TURNS * 2
            for role, content in list(user.history)[-keep_turns:]:
                truncated = content[:500] + "..." if len(content) > 500 else content
                history_messages.append({
                    "role": "user" if role == "user" else "assistant",
                    "content": truncated
                })
            
            messages.extend(history_messages)
        
        messages.append({"role": "user", "content": user_text})
        
        return messages

message_builder = MessageBuilder()

async def stream_response(
    context: ContextTypes.DEFAULT_TYPE,
    chat_id: int,
    model: str,
    messages: List[Dict[str, str]],
    max_tokens: int,
    temperature: float = 0.7
):
    msg = await context.bot.send_message(chat_id, "💭 Đang suy nghĩ...")
    
    full_response = ""
    chunk_buffer = ""
    update_counter = 0
    
    async def update_message():
        nonlocal chunk_buffer
        try:
            if chunk_buffer:
                await msg.edit_text(chunk_buffer[:4096])
        except:
            pass
    
    try:
        stream = await asyncio.to_thread(
            ai_client.stream_complete,
            model, messages, max_tokens, temperature
        )
        
        for chunk in stream:
            full_response += chunk
            chunk_buffer += chunk
            update_counter += 1
            
            if update_counter % 5 == 0 and len(chunk_buffer) > 100:
                await update_message()
        
        if full_response:
            chunks = text_processor.chunk_text(full_response)
            
            if len(chunks) == 1:
                await msg.edit_text(chunks[0])
            else:
                await msg.delete()
                for i, chunk in enumerate(chunks):
                    await context.bot.send_message(
                        chat_id,
                        f"📄 Phần {i+1}/{len(chunks)}:\n\n{chunk}"
                    )
        
        return full_response
        
    except Exception as e:
        await msg.edit_text(f"❌ Lỗi: {str(e)[:200]}")
        return None

# Command Handlers
async def cmd_help(update: Update, context: ContextTypes.DEFAULT_TYPE):
    help_text = """
🇻🇳 **LINH AI - TRỢ LÝ VIỆT NAM**

📝 **Lệnh cơ bản:**
• /help - Hướng dẫn sử dụng
• /start - Khởi động bot
• /clear - Xóa lịch sử chat
• /stats - Thống kê sử dụng

💻 **Tính năng AI:**
• /code <yêu cầu> - Viết code chuyên nghiệp
• /img <mô tả> - Tạo ảnh AI (Gemini)
• /vision - Phân tích ảnh đã gửi

🇻🇳 **Việt Nam:**
• /weather <city> - Thời tiết các tỉnh thành
• /news - Tin tức mới nhất
• /exchange - Tỷ giá ngoại tệ
• /time - Giờ Việt Nam
• /translate <text> - Dịch Anh-Việt

📸 **Phân tích ảnh:**
• Gửi ảnh → Bot phân tích tự động
• Hỏi chi tiết về nội dung ảnh
• So sánh nhiều ảnh cùng lúc
• Nhận diện văn hóa Việt Nam

📄 **Xử lý File:**
• Hỗ trợ: PDF, Word, Excel, PowerPoint
• Text, Code, JSON, CSV, HTML
• /cancelfile - Thoát chế độ file
• /sendfile - Tải kết quả về

💡 **Mẹo hay:**
• Chat tiếng Việt tự nhiên
• "img: <mô tả>" để tạo ảnh nhanh
• Hỏi về lịch sử, ẩm thực, du lịch VN

⚙️ **AI Models:**
• Chat: Claude-3.5 (Vercel)
• Vision: Gemini-1.5-Flash
• Image: Gemini-2.0-Flash

👨‍💻 Dev: @cucodoivandep
🌐 Made in Vietnam with ❤️
    """
    
    await context.bot.send_message(
        update.effective_chat.id,
        help_text,
        parse_mode=ParseMode.MARKDOWN
    )

async def cmd_start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    welcome_text = """
🇻🇳 **Xin chào! Mình là Linh - AI Assistant Việt Nam**

🎯 **Mình có thể giúp gì cho bạn:**
• 📸 Phân tích hình ảnh, nhận diện đối tượng
• 💬 Trò chuyện về mọi chủ đề
• 💻 Viết code, debug, tối ưu
• 📚 Tư vấn học tập, công việc
• 🇻🇳 Thông tin về Việt Nam

💡 **Thử ngay:**
• Gửi ảnh để phân tích
• Hỏi về văn hóa, lịch sử Việt Nam
• /help để xem đầy đủ tính năng

Chúc bạn một ngày tốt lành! 🌺
    """
    
    keyboard = InlineKeyboardMarkup([
        [
            InlineKeyboardButton("📖 Hướng dẫn", callback_data="show_help"),
            InlineKeyboardButton("🇻🇳 Về Việt Nam", callback_data="about_vietnam")
        ],
        [
            InlineKeyboardButton("📸 Cách dùng ảnh", callback_data="image_guide"),
            InlineKeyboardButton("💬 Chat ngay", callback_data="start_chat")
        ]
    ])
    
    await context.bot.send_message(
        update.effective_chat.id,
        welcome_text,
        parse_mode=ParseMode.MARKDOWN,
        reply_markup=keyboard
    )

async def cmd_vision(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Analyze pending images"""
    chat_id = update.effective_chat.id
    user = bot_state.get_user(chat_id)
    
    if not user.pending_images:
        await context.bot.send_message(
            chat_id,
            "📷 Chưa có ảnh nào.\nGửi ảnh để phân tích nhé!"
        )
        return
    
    if not config.GEMINI_API_KEY:
        await context.bot.send_message(
            chat_id,
            "❌ Cần GEMINI_API_KEY để phân tích ảnh.\nLiên hệ @cucodoivandep"
        )
        return
    
    await context.bot.send_message(
        chat_id,
        f"🔍 Đang phân tích {len(user.pending_images)} ảnh..."
    )
    
    await context.bot.send_chat_action(chat_id=chat_id, action=ChatAction.TYPING)
    
    prompt = "Phân tích chi tiết hình ảnh: nội dung, đối tượng, màu sắc, bố cục. Nếu có yếu tố Việt Nam (người, cảnh, món ăn...) hãy mô tả kỹ."
    
    analysis = await gemini_handler.analyze_multiple_images(
        user.pending_images,
        prompt
    )
    
    chunks = text_processor.chunk_text(analysis)
    for chunk in chunks:
        await context.bot.send_message(
            chat_id, 
            chunk,
            parse_mode=ParseMode.MARKDOWN
        )
    
    user.history.append(("user", "Phân tích ảnh"))
    user.history.append(("assistant", analysis[:500]))

async def cmd_img(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Generate image from text"""
    chat_id = update.effective_chat.id
    prompt = " ".join(context.args).strip()
    
    if not prompt:
        examples = [
            "phong cảnh vịnh Hạ Long",
            "phở bò Hà Nội",
            "áo dài Việt Nam",
            "chợ nổi Cái Răng",
            "ruộng bậc thang Sapa"
        ]
        await context.bot.send_message(
            chat_id,
            f"📝 **Cú pháp:** /img <mô tả>\n\n"
            f"**Ví dụ:**\n" + "\n".join([f"• /img {ex}" for ex in examples])
        )
        return
    
    if not config.GEMINI_API_KEY:
        await context.bot.send_message(
            chat_id,
            "❌ Cần GEMINI_API_KEY để tạo ảnh.\nLiên hệ @cucodoivandep"
        )
        return
    
    await context.bot.send_chat_action(chat_id=chat_id, action=ChatAction.UPLOAD_PHOTO)
    
    # Enhance prompt for Vietnamese context
    enhanced_prompt = f"{prompt}, high quality, detailed, beautiful"
    
    status_msg = await context.bot.send_message(
        chat_id,
        f"🎨 Đang tạo ảnh: _{prompt}_",
        parse_mode=ParseMode.MARKDOWN
    )
    
    # Try to generate with Gemini
    image_data = await gemini_handler.generate_image(enhanced_prompt)
    
    if image_data:
        await status_msg.delete()
        await context.bot.send_photo(
            chat_id,
            photo=io.BytesIO(image_data),
            caption=f"🎨 {prompt}"
        )
    else:
        await status_msg.edit_text(
            "❌ Không thể tạo ảnh. Gemini có thể đang bận hoặc prompt không phù hợp.\n"
            "💡 Thử lại với mô tả khác nhé!"
        )
    
    user = bot_state.get_user(chat_id)
    user.history.append(("user", f"/img {prompt}"))
    user.history.append(("assistant", f"Tạo ảnh: {prompt[:50]}"))

async def cmd_clear(update: Update, context: ContextTypes.DEFAULT_TYPE):
    chat_id = update.effective_chat.id
    user = bot_state.get_user(chat_id)
    
    # Clear everything
    user.history.clear()
    user.file_mode = False
    user.image_mode = False
    user.pending_file = None
    user.pending_images.clear()
    user.last_result = ""
    
    await context.bot.send_message(
        chat_id,
        "✅ Đã xóa:\n"
        "• Lịch sử chat\n"
        "• Ảnh đã gửi\n"
        "• File đã gửi\n"
        "• Kết quả lưu"
    )

async def cmd_stats(update: Update, context: ContextTypes.DEFAULT_TYPE):
    chat_id = update.effective_chat.id
    user = bot_state.get_user(chat_id)
    
    try:
        import pytz
        vn_tz = pytz.timezone('Asia/Ho_Chi_Minh')
        vn_time = datetime.datetime.now(vn_tz)
        time_str = vn_time.strftime('%H:%M:%S %d/%m/%Y')
    except:
        vn_time = datetime.datetime.utcnow() + datetime.timedelta(hours=7)
        time_str = vn_time.strftime('%H:%M:%S %d/%m/%Y')
    
    stats_text = f"""
📊 **Thống kê sử dụng**

👤 **User:** {update.effective_user.first_name}
🆔 **ID:** `{chat_id}`
💬 **Lịch sử:** {len(user.history)} tin
📸 **Ảnh lưu:** {len(user.pending_images)}
📁 **File mode:** {'Bật' if user.file_mode else 'Tắt'}
🌍 **Vị trí:** {user.location or 'Hà Nội'}

⚙️ **AI Models:**
• Chat: Claude-3.5-Haiku
• Code: Claude-3.5-Sonnet
• Vision: Gemini-1.5-Flash
• Image: Gemini-2.0-Flash

🕐 **Giờ VN:** {time_str}

💡 _Dùng /clear để xóa dữ liệu_
    """
    
    await context.bot.send_message(
        chat_id,
        stats_text,
        parse_mode=ParseMode.MARKDOWN
    )

async def cmd_weather(update: Update, context: ContextTypes.DEFAULT_TYPE):
    chat_id = update.effective_chat.id
    
    if not context.args:
        cities = ", ".join(VIETNAM_CITIES.keys())
        await context.bot.send_message(
            chat_id,
            f"☀️ **Xem thời tiết**\n\n"
            f"Cú pháp: /weather <tên thành phố>\n\n"
            f"Các thành phố: {cities}"
        )
        return
    
    city = context.args[0].lower()
    weather_info = await vietnam_services.get_weather(city)
    
    await context.bot.send_message(
        chat_id,
        weather_info,
        parse_mode=ParseMode.MARKDOWN
    )

async def cmd_code(update: Update, context: ContextTypes.DEFAULT_TYPE):
    chat_id = update.effective_chat.id
    request = " ".join(context.args).strip()
    
    if not request:
        examples = [
            "hàm tính giai thừa Python",
            "validate email JavaScript", 
            "quicksort C++",
            "REST API với Flask",
            "React component với hooks"
        ]
        await context.bot.send_message(
            chat_id,
            f"💻 **Viết code**\n\n"
            f"Cú pháp: /code <yêu cầu>\n\n"
            f"**Ví dụ:**\n" + "\n".join([f"• /code {ex}" for ex in examples])
        )
        return
    
    await context.bot.send_chat_action(chat_id=chat_id, action=ChatAction.TYPING)
    
    messages = message_builder.build_messages(
        chat_id, 
        request,
        context_type="code",
        include_history=False
    )
    
    result = await stream_response(
        context,
        chat_id,
        config.CODE_MODEL,
        messages,
        config.MAX_TOKENS_CODE,
        temperature=0.3
    )
    
    if result:
        user = bot_state.get_user(chat_id)
        user.last_result = result
        user.history.append(("user", f"/code {request}"))
        user.history.append(("assistant", result[:500]))

async def cmd_translate(update: Update, context: ContextTypes.DEFAULT_TYPE):
    chat_id = update.effective_chat.id
    text = " ".join(context.args).strip()
    
    if not text:
        await context.bot.send_message(
            chat_id,
            "🔤 **Dịch Anh-Việt**\n\n"
            "Cú pháp: /translate <text tiếng Anh>\n\n"
            "Ví dụ: /translate Hello world"
        )
        return
    
    messages = [
        {
            "role": "system", 
            "content": "Bạn là chuyên gia dịch thuật. Dịch chính xác, tự nhiên sang tiếng Việt. Chỉ trả về bản dịch, không giải thích."
        },
        {
            "role": "user",
            "content": f"Dịch sang tiếng Việt:\n{text}"
        }
    ]
    
    result = await ai_client.complete(
        config.CHAT_MODEL,
        messages,
        500,
        temperature=0.3
    )
    
    await context.bot.send_message(
        chat_id,
        f"🔤 **Bản gốc:**\n{text}\n\n"
        f"🇻🇳 **Bản dịch:**\n{result}",
        parse_mode=ParseMode.MARKDOWN
    )

async def cmd_news(update: Update, context: ContextTypes.DEFAULT_TYPE):
    chat_id = update.effective_chat.id
    
    await context.bot.send_chat_action(chat_id=chat_id, action=ChatAction.TYPING)
    news = await vietnam_services.get_news()
    
    await context.bot.send_message(
        chat_id,
        news,
        parse_mode=ParseMode.MARKDOWN
    )

async def cmd_exchange(update: Update, context: ContextTypes.DEFAULT_TYPE):
    chat_id = update.effective_chat.id
    
    rate = await vietnam_services.get_exchange_rate()
    
    await context.bot.send_message(
        chat_id,
        rate,
        parse_mode=ParseMode.MARKDOWN
    )

async def cmd_time(update: Update, context: ContextTypes.DEFAULT_TYPE):
    chat_id = update.effective_chat.id
    
    time_info = vietnam_services.get_vietnam_time()
    
    await context.bot.send_message(
        chat_id,
        time_info,
        parse_mode=ParseMode.MARKDOWN
    )

async def cmd_sendfile(update: Update, context: ContextTypes.DEFAULT_TYPE):
    chat_id = update.effective_chat.id
    user = bot_state.get_user(chat_id)
    
    if not user.last_result:
        await context.bot.send_message(
            chat_id,
            "❌ Không có kết quả để gửi.\n"
            "💡 Chat hoặc dùng /code trước"
        )
        return
    
    file_bytes = user.last_result.encode('utf-8')
    file_io = io.BytesIO(file_bytes)
    file_io.name = f"result_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
    
    await context.bot.send_document(
        chat_id,
        document=file_io,
        caption="📄 Kết quả xử lý"
    )

async def cmd_cancelfile(update: Update, context: ContextTypes.DEFAULT_TYPE):
    chat_id = update.effective_chat.id
    user = bot_state.get_user(chat_id)
    
    user.file_mode = False
    user.image_mode = False
    user.pending_file = None
    
    await context.bot.send_message(
        chat_id,
        "✅ Đã thoát chế độ file/image"
    )

# Message Handlers
async def on_photo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle photo messages"""
    chat_id = update.effective_chat.id
    user = bot_state.get_user(chat_id)
    
    if not config.GEMINI_API_KEY:
        await context.bot.send_message(
            chat_id,
            "❌ Cần GEMINI_API_KEY để phân tích ảnh.\n"
            "Liên hệ @cucodoivandep để được hỗ trợ."
        )
        return
    
    # Get photo
    photo = update.message.photo[-1]
    caption = update.message.caption or ""
    
    await context.bot.send_chat_action(chat_id=chat_id, action=ChatAction.TYPING)
    
    try:
        # Download photo
        file = await context.bot.get_file(photo.file_id)
        file_data = await file.download_as_bytearray()
        
        # Store image
        img_data = ImageData(
            file_id=photo.file_id,
            file_data=bytes(file_data),
            caption=caption
        )
        user.pending_images.append(img_data)
        user.image_mode = True
        
        # Quick analysis
        prompt = (
            "Phân tích chi tiết ảnh này bằng tiếng Việt:\n"
            "1. Mô tả nội dung chính\n"
            "2. Nhận diện đối tượng, con người\n"
            "3. Màu sắc, bố cục\n"
            "4. Nếu có yếu tố Việt Nam, hãy nhấn mạnh"
        )
        
        if caption:
            prompt += f"\n\nThông tin thêm: {caption}"
        
        analysis = await gemini_handler.analyze_image(
            img_data.file_data,
            prompt
        )
        
        if analysis:
            response = (
                f"📸 **Ảnh #{len(user.pending_images)}**\n\n"
                f"{analysis}\n\n"
                f"💬 Bạn có thể:\n"
                f"• Hỏi chi tiết về ảnh\n"
                f"• Gửi thêm ảnh để so sánh\n"
                f"• /vision xem lại tất cả\n"
                f"• /clear xóa ảnh"
            )
        else:
            response = (
                f"📸 Đã nhận ảnh #{len(user.pending_images)}\n"
                f"❌ Không thể phân tích ngay.\n"
                f"Hãy hỏi cụ thể về ảnh nhé!"
            )
        
        await context.bot.send_message(
            chat_id, 
            response,
            parse_mode=ParseMode.MARKDOWN
        )
        
        # Save to history
        if analysis:
            user.history.append(("user", f"[Gửi ảnh] {caption[:50] if caption else ''}"))
            user.history.append(("assistant", analysis[:200]))
        
    except Exception as e:
        await context.bot.send_message(
            chat_id,
            f"❌ Lỗi xử lý ảnh: {str(e)[:100]}"
        )

async def on_file(update: Update, context: ContextTypes.DEFAULT_TYPE):
    chat_id = update.effective_chat.id
    document = update.message.document
    
    if not document:
        return
    
    if document.file_size > 20 * 1024 * 1024:
        await context.bot.send_message(
            chat_id,
            "❌ File quá lớn (max 20MB)"
        )
        return
    
    await context.bot.send_chat_action(chat_id=chat_id, action=ChatAction.TYPING)
    
    try:
        file = await context.bot.get_file(document.file_id)
        file_data = await file.download_as_bytearray()
        
        content, file_type = await file_processor.process_file(
            document.file_name or "unknown",
            bytes(file_data)
        )
        
        if file_type == "image":
            # Handle as image
            user = bot_state.get_user(chat_id)
            img_data = ImageData(
                file_id=document.file_id,
                file_data=bytes(file_data),
                caption=document.file_name
            )
            user.pending_images.append(img_data)
            user.image_mode = True
            
            await context.bot.send_message(
                chat_id,
                f"📷 Đã nhận ảnh từ file\n"
                f"• Hỏi về nội dung ảnh\n"
                f"• /vision để phân tích"
            )
            return
        
        if not content:
            await context.bot.send_message(chat_id, file_type)
            return
        
        user = bot_state.get_user(chat_id)
        user.file_mode = True
        user.pending_file = {
            "name": document.file_name,
            "content": content[:config.CHUNK_CHARS],
            "type": file_type
        }
        
        preview = content[:500] + "..." if len(content) > 500 else content
        
        await context.bot.send_message(
            chat_id,
            f"✅ **File đã nhận**\n\n"
            f"📄 Tên: {document.file_name}\n"
            f"📊 Loại: {file_type}\n"
            f"📝 Kích thước: {len(content):,} ký tự\n\n"
            f"**Xem trước:**\n```\n{preview}\n```\n\n"
            f"💬 Hỏi về file hoặc /cancelfile",
            parse_mode=ParseMode.MARKDOWN
        )
        
    except Exception as e:
        await context.bot.send_message(
            chat_id,
            f"❌ Lỗi: {str(e)[:100]}"
        )

async def on_text(update: Update, context: ContextTypes.DEFAULT_TYPE):
    chat_id = update.effective_chat.id
    message_text = (update.message.text or "").strip()
    
    if not message_text:
        return
    
    # Quick image generation
    if message_text.lower().startswith("img:"):
        prompt = message_text[4:].strip()
        context.args = prompt.split()
        await cmd_img(update, context)
        return
    
    user = bot_state.get_user(chat_id)
    
    await context.bot.send_chat_action(chat_id=chat_id, action=ChatAction.TYPING)
    
    # Handle image mode
    if user.image_mode and user.pending_images:
        if not config.GEMINI_API_KEY:
            await context.bot.send_message(
                chat_id,
                "❌ Cần GEMINI_API_KEY để phân tích ảnh"
            )
            return
        
        # Analyze images with question
        await context.bot.send_message(
            chat_id,
            f"🔍 Đang phân tích {len(user.pending_images)} ảnh..."
        )
        
        results = []
        for i, img_data in enumerate(user.pending_images, 1):
            prompt = f"Với câu hỏi: '{message_text}'\nHãy phân tích ảnh và trả lời bằng tiếng Việt."
            
            analysis = await gemini_handler.analyze_image(
                img_data.file_data,
                prompt
            )
            
            if analysis:
                results.append(f"**📸 Ảnh {i}:**\n{analysis}")
        
        if results:
            full_response = "\n\n".join(results)
            
            # Also get context from text AI
            image_context = f"Người dùng đã gửi {len(user.pending_images)} ảnh và hỏi: {message_text}\nPhân tích ảnh cho thấy: {full_response[:500]}"
            
            messages = message_builder.build_messages(
                chat_id,
                image_context,
                context_type="image_context",
                include_history=True
            )
            
            ai_response = await ai_client.complete(
                config.CHAT_MODEL,
                messages,
                config.MAX_TOKENS,
                temperature=0.7
            )
            
            # Combine results
            final_response = full_response
            if ai_response and not ai_response.startswith("❌"):
                final_response += f"\n\n💭 **Nhận xét thêm:**\n{ai_response}"
            
            chunks = text_processor.chunk_text(final_response)
            for chunk in chunks:
                await context.bot.send_message(
                    chat_id,
                    chunk,
                    parse_mode=ParseMode.MARKDOWN
                )
            
            user.last_result = final_response
            user.history.append(("user", message_text))
            user.history.append(("assistant", final_response[:500]))
        else:
            await context.bot.send_message(
                chat_id,
                "❌ Không thể phân tích ảnh"
            )
        
        return
    
    # Handle file mode
    if user.file_mode and user.pending_file:
        file_content = user.pending_file["content"]
        file_name = user.pending_file["name"]
        
        prompt = (
            f"File: {file_name}\n"
            f"Nội dung:\n{file_content[:10000]}\n\n"
            f"Câu hỏi: {message_text}"
        )
        
        messages = message_builder.build_messages(
            chat_id,
            prompt,
            context_type="file",
            include_history=False
        )
        
        result = await stream_response(
            context,
            chat_id,
            config.FILE_MODEL,
            messages,
            config.FILE_OUTPUT_TOKENS,
            temperature=0.5
        )
    else:
        # Normal chat
        messages = message_builder.build_messages(
            chat_id,
            message_text,
            context_type="chat",
            include_history=True
        )
        
        result = await stream_response(
            context,
            chat_id,
            config.CHAT_MODEL,
            messages,
            config.MAX_TOKENS,
            temperature=0.7
        )
    
    if result:
        user.last_result = result
        user.history.append(("user", message_text[:500]))
        user.history.append(("assistant", result[:500]))

async def on_callback_query(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    
    if query.data == "show_help":
        await cmd_help(update, context)
    
    elif query.data == "about_vietnam":
        vietnam_info = """
🇻🇳 **VIỆT NAM - ĐẤT NƯỚC CON NGƯỜI**

🏛 **Lịch sử vẻ vang:**
• 4000 năm văn hiến
• 18 đời vua Hùng dựng nước
• Đánh bại Mông-Nguyên 3 lần
• Độc lập thống nhất năm 1975

🌏 **Địa lý tươi đẹp:**
• Diện tích: 331,212 km²
• Dân số: ~98 triệu người
• 3,260 km bờ biển
• 2 đồng bằng phì nhiêu

🎭 **Văn hóa đa dạng:**
• 54 dân tộc anh em
• 8 Di sản UNESCO
• Tết Nguyên Đán độc đáo
• Ẩm thực phong phú

🏆 **Thành tựu hiện đại:**
• Top 20 nền kinh tế lớn nhất
• Xuất khẩu gạo số 2 thế giới
• Du lịch phát triển mạnh
• Công nghệ số bùng nổ

💪 **Việt Nam - Đất nước anh hùng!**
🌟 **Tiềm năng - Khát vọng - Vươn cao!**
        """
        await context.bot.send_message(
            query.message.chat_id,
            vietnam_info,
            parse_mode=ParseMode.MARKDOWN
        )
    
    elif query.data == "image_guide":
        guide = """
📸 **HƯỚNG DẪN DÙNG ẢNH**

**Cách gửi:**
1️⃣ Gửi 1 hoặc nhiều ảnh
2️⃣ Bot phân tích tự động
3️⃣ Hỏi chi tiết về ảnh

**Bot có thể:**
✅ Nhận diện đối tượng
✅ Đọc chữ trong ảnh (OCR)
✅ Phân tích màu sắc, bố cục
✅ Nhận diện món ăn VN
✅ Nhận diện địa danh VN
✅ So sánh nhiều ảnh

**Câu hỏi mẫu:**
• "Đây là món gì?"
• "Có bao nhiêu người?"
• "Đây là ở đâu?"
• "Dịch chữ trong ảnh"
• "So sánh 2 ảnh này"

**Lệnh:**
• /vision - Xem lại phân tích
• /clear - Xóa ảnh đã gửi

💡 Gửi ảnh rõ nét để có kết quả tốt nhất!
        """
        await context.bot.send_message(
            query.message.chat_id,
            guide,
            parse_mode=ParseMode.MARKDOWN
        )
    
    elif query.data == "start_chat":
        await context.bot.send_message(
            query.message.chat_id,
            "💬 Sẵn sàng! Hãy chat với mình.\n"
            "Bạn có thể hỏi về Việt Nam, gửi ảnh để phân tích, hoặc bất cứ điều gì! 😊"
        )

def main():
    if not config.BOT_TOKEN:
        print("❌ Thiếu BOT_TOKEN")
        return
    
    print("=" * 50)
    print("🇻🇳 LINH BOT - AI Assistant Việt Nam")
    print("=" * 50)
    
    if not config.VERCEL_API_KEY:
        print("⚠️  Thiếu VERCEL_API_KEY - Chat/Code bị hạn chế")
    else:
        print("✅ Vercel API: OK")
    
    if not config.GEMINI_API_KEY:
        print("⚠️  Thiếu GEMINI_API_KEY - Vision/Image bị hạn chế")
    else:
        print("✅ Gemini API: OK")
    
    print("=" * 50)
    
    app = ApplicationBuilder().token(config.BOT_TOKEN).build()
    
    # Commands
    app.add_handler(CommandHandler("start", cmd_start))
    app.add_handler(CommandHandler("help", cmd_help))
    app.add_handler(CommandHandler("clear", cmd_clear))
    app.add_handler(CommandHandler("stats", cmd_stats))
    app.add_handler(CommandHandler("img", cmd_img))
    app.add_handler(CommandHandler("code", cmd_code))
    app.add_handler(CommandHandler("vision", cmd_vision))
    app.add_handler(CommandHandler("weather", cmd_weather))
    app.add_handler(CommandHandler("news", cmd_news))
    app.add_handler(CommandHandler("exchange", cmd_exchange))
    app.add_handler(CommandHandler("time", cmd_time))
    app.add_handler(CommandHandler("translate", cmd_translate))
    app.add_handler(CommandHandler("sendfile", cmd_sendfile))
    app.add_handler(CommandHandler("cancelfile", cmd_cancelfile))
    
    # Callbacks
    app.add_handler(CallbackQueryHandler(on_callback_query))
    
    # Messages
    app.add_handler(MessageHandler(filters.PHOTO, on_photo))
    app.add_handler(MessageHandler(filters.Document.ALL, on_file))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, on_text))
    
    print("🚀 Bot đang chạy...")
    print("📸 Vision: Gemini-1.5-Flash")
    print("🎨 Image Gen: Gemini-2.0-Flash")
    print("💬 Chat: Claude-3.5 via Vercel")
    print("👨‍💻 Dev: @cucodoivandep")
    print("=" * 50)
    
    app.run_polling(drop_pending_updates=True)

if __name__ == "__main__":
    main()
